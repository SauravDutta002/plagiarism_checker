import os
import fitz  
import docx  
import pytesseract  
from PIL import Image  
from flask import Flask, render_template, request
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import requests  
from pptx import Presentation  

app = Flask(__name__)

SERP_API_KEY = "c827578ca2241ad2e587c7dff8b97bd7eff06dd192d297e72e8a86fb4e9c7fd9"  

def search_online_sources(query):
    """Fetch matching sources from the internet using SerpAPI."""
    url = "https://serpapi.com/search"
    params = {
        "engine": "google",
        "q": query,
        "api_key": SERP_API_KEY,
        "num": 5  # Get top 5 results
    }
    
    response = requests.get(url, params=params)
    data = response.json()
    
    sources = []
    if "organic_results" in data:
        for result in data["organic_results"]:
            sources.append({"text": result["title"], "url": result["link"]})
    
    return sources

def check_plagiarism(input_text):
    """Checks plagiarism of the given text by searching online sources."""
    input_text = input_text.lower().strip()
    
    if not input_text:
        return 0, []

    sources = search_online_sources(input_text)

    source_texts = [source["text"] for source in sources]
    
    if not source_texts:
        return 0, []  
    
    vectorizer = TfidfVectorizer().fit_transform([input_text] + source_texts)
    similarity_matrix = cosine_similarity(vectorizer[0], vectorizer[1:])  


    max_similarity = np.max(similarity_matrix) * 100  

    return round(max_similarity, 2), sources

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text("text") + "\n"
    return text.strip()

def extract_text_from_docx(docx_path):
    """Extract text from a DOCX file."""
    doc = docx.Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs]).strip()

def extract_text_from_pptx(pptx_path):
    """Extract text from a PPTX file."""
    prs = Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text).strip()

def extract_text_from_image(image_path):
    """Extract text from an image using OCR."""
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text.strip()

@app.route("/", methods=["GET", "POST"])
def index():
    plagiarism_percentage = None
    sources = []
    text = None

    if request.method == "POST":
        if "text" in request.form and request.form["text"].strip():
            text = request.form["text"].strip()
            plagiarism_percentage, sources = check_plagiarism(text)

        elif "file" in request.files:
            uploaded_file = request.files["file"]

            if uploaded_file and uploaded_file.filename:
                filename = uploaded_file.filename
                file_ext = filename.split(".")[-1].lower()
                file_path = os.path.join("uploads", filename)

                uploaded_file.save(file_path)

                if file_ext == "pdf":
                    text = extract_text_from_pdf(file_path)
                elif file_ext == "docx":
                    text = extract_text_from_docx(file_path)
                elif file_ext == "pptx":
                    text = extract_text_from_pptx(file_path)
                elif file_ext in ["png", "jpg", "jpeg"]:
                    text = extract_text_from_image(file_path)
                
                if text:
                    plagiarism_percentage, sources = check_plagiarism(text)
                
                os.remove(file_path)  # Clean up

        return render_template("result.html", text=text, plagiarism_percentage=plagiarism_percentage, sources=sources)

    return render_template("index.html")

if __name__ == "__main__":
    if not os.path.exists("uploads"):
        os.makedirs("uploads")
    app.run(debug=True)
